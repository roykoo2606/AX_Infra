#!/usr/bin/env python3
"""Convert office/hancom/pdf documents into agent-readable Markdown.

Usage: parse-docs.py <src-dir> <out-dir>

Mirrors the source tree into <out-dir>, writing one .md per source file.
Tables are preserved as Markdown tables wherever the source format allows.
Supported: .hwp .hwpx .pptx .xlsx .docx .pdf .csv .html/.htm .zip
Other types get a stub noting the file exists but carries no extractable text.
"""
import csv
import io
import json
import subprocess
import sys
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path

MAX_ROWS = 200        # per sheet / table
MAX_CHARS = 400_000   # per output document


def md_table(rows):
    """Render a list of row-lists as a Markdown table."""
    rows = [[(c if c is not None else "").replace("\n", " ").replace("|", "\\|").strip()
             for c in r] for r in rows]
    rows = [r for r in rows if any(c for c in r)]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    head, body = rows[0], rows[1:]
    out = ["| " + " | ".join(head) + " |",
           "|" + "|".join(["---"] * width) + "|"]
    out += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(out)


# --- Hancom -----------------------------------------------------------------

def local(el):
    t = el.tag
    return t.rsplit("}", 1)[-1] if isinstance(t, str) and "}" in t else t


def hwpx_container(el, out, depth=0):
    """Walk an HWPX container in document order, emitting paragraphs and tables."""
    import lxml.etree as ET  # noqa: F401

    def cell_text(tc):
        buf = []
        hwpx_container(tc, buf, depth + 1)
        return " ".join(b for b in buf if b and not b.startswith("|"))

    def render_tbl(tbl):
        rows = []
        for tr in tbl.iter():
            if local(tr) != "tr":
                continue
            rows.append([cell_text(tc) for tc in tr if local(tc) == "tc"])
        return md_table(rows)

    buf = []

    def flush():
        if buf:
            text = "".join(buf).strip()
            if text:
                out.append(text)
            buf.clear()

    def rec(node):
        for ch in node:
            name = local(ch)
            if name == "tbl":
                flush()
                if depth < 4:
                    t = render_tbl(ch)
                    if t:
                        out.append(t)
            elif name == "t":
                buf.append("".join(ch.itertext()))
            elif name == "lineBreak":
                buf.append("\n")
            else:
                rec(ch)

    for child in el:
        if local(child) == "p":
            rec(child)
            flush()
        else:
            rec(child)
            flush()


def parse_hwpx(path):
    import lxml.etree as ET
    out = []
    with zipfile.ZipFile(path) as z:
        sections = sorted(n for n in z.namelist()
                          if n.startswith("Contents/section") and n.endswith(".xml"))
        for name in sections:
            root = ET.fromstring(z.read(name))
            hwpx_container(root, out)
    return "\n\n".join(out)


PYHWP_PY = Path.home() / ".local/share/uv/tools/pyhwp/bin/python"
RESCUE = Path(__file__).resolve().parent / "hwp5_rescue.py"


def hwp_rescue(path):
    """Last resort for HWP files pyhwp rejects (e.g. missing summary stream)."""
    if not (PYHWP_PY.exists() and RESCUE.exists()):
        return ""
    r = subprocess.run([str(PYHWP_PY), str(RESCUE), str(path)],
                       capture_output=True, timeout=180)
    return r.stdout.decode("utf-8", "replace")


def parse_hwp(path):
    """hwp5html preserves tables; fall back to hwp5txt, then to the rescue reader."""
    import lxml.html as LH
    with tempfile.TemporaryDirectory() as tmp:
        try:
            subprocess.run(["hwp5html", "--output", tmp, str(path)],
                           check=True, capture_output=True, timeout=120)
            doc = LH.parse(str(Path(tmp) / "index.xhtml"))
        except Exception:
            r = subprocess.run(["hwp5txt", str(path)],
                               capture_output=True, timeout=120)
            text = r.stdout.decode("utf-8", "replace").strip()
            return text or hwp_rescue(path)

        out = []
        body = doc.getroot().body if doc.getroot().body is not None else doc.getroot()

        def walk(el):
            for ch in el:
                tag = ch.tag if isinstance(ch.tag, str) else ""
                if tag == "table":
                    rows = [[" ".join("".join(td.itertext()).split())
                             for td in tr.xpath("./td|./th")]
                            for tr in ch.xpath(".//tr")]
                    t = md_table(rows)
                    if t:
                        out.append(t)
                elif tag in ("p", "div", "h1", "h2", "h3", "li"):
                    if ch.xpath(".//table"):
                        walk(ch)
                    else:
                        text = " ".join("".join(ch.itertext()).split())
                        if text:
                            out.append(text)
                else:
                    walk(ch)

        walk(body)
        return "\n\n".join(out)


# --- Office -----------------------------------------------------------------

def parse_pptx(path):
    from pptx import Presentation
    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"## 슬라이드 {i}"]
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                t = md_table(rows)
                if t:
                    parts.append(t)
            elif shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if text.strip():
                    parts.append(text.strip())
        try:
            if slide.has_notes_slide:
                note = slide.notes_slide.notes_text_frame.text.strip()
                if note:
                    parts.append(f"**[발표자 노트]**\n{note}")
        except Exception:
            pass
        out.append("\n\n".join(parts))
    return "\n\n".join(out)


def parse_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    out = []
    for ws in wb.worksheets:
        rows = []
        for r in ws.iter_rows(max_row=MAX_ROWS, values_only=True):
            rows.append(["" if v is None else str(v) for v in r])
        while rows and not any(c.strip() for c in rows[-1]):
            rows.pop()
        if not rows:
            continue
        header = f"## 시트: {ws.title}  (총 {ws.max_row}행 × {ws.max_column}열)"
        note = f"\n> 앞 {MAX_ROWS}행만 표시" if (ws.max_row or 0) > MAX_ROWS else ""
        out.append(f"{header}{note}\n\n{md_table(rows)}")
    wb.close()
    return "\n\n".join(out)


def parse_docx(path):
    import docx
    d = docx.Document(str(path))
    out = []
    for p in d.paragraphs:
        if p.text.strip():
            out.append(p.text.strip())
    for tbl in d.tables:
        rows = [[c.text for c in row.cells] for row in tbl.rows]
        t = md_table(rows)
        if t:
            out.append(t)
    return "\n\n".join(out)


# --- PDF / plain ------------------------------------------------------------

def parse_pdf(path):
    r = subprocess.run(["pdftotext", "-layout", "-enc", "UTF-8", str(path), "-"],
                       capture_output=True, timeout=300)
    text = r.stdout.decode("utf-8", "replace").strip()
    if len(text) < 40:  # likely a scan with no text layer
        try:
            import fitz
            doc = fitz.open(str(path))
            pages, imgs = doc.page_count, sum(len(p.get_images()) for p in doc)
            doc.close()
            return (f"_(텍스트 레이어 없음 — 스캔 이미지 PDF로 판단. "
                    f"{pages}페이지, 이미지 {imgs}개. 내용 확인은 원본 열람 필요)_")
        except Exception:
            return "_(텍스트 추출 불가)_"
    return text


def parse_csv(path):
    raw = path.read_bytes()
    for enc in ("utf-8-sig", "utf-8", "cp949", "euc-kr"):
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    else:
        return "_(인코딩 판별 실패)_"
    rows = list(csv.reader(io.StringIO(text)))
    if not rows:
        return "_(빈 파일)_"
    head, data = rows[0], rows[1:]
    out = [f"- 인코딩: `{enc}`", f"- 총 {len(data):,}행 × {len(head)}열",
           f"- 컬럼: {', '.join(head)}", "",
           f"### 앞 {min(MAX_ROWS, len(data))}행", "",
           md_table([head] + data[:MAX_ROWS])]
    # per-column distinct-value summary, useful for schema understanding
    summary = []
    for i, col in enumerate(head):
        vals = [r[i] for r in data if i < len(r) and r[i].strip()]
        uniq = sorted(set(vals))
        sample = ", ".join(uniq[:8]) + (" …" if len(uniq) > 8 else "")
        summary.append([col, str(len(vals)), str(len(uniq)), sample])
    out += ["", "### 컬럼 요약", "",
            md_table([["컬럼", "값있음", "고유값", "값 예시"]] + summary)]
    return "\n".join(out)


def parse_html(path):
    import lxml.html as LH
    doc = LH.fromstring(path.read_bytes())
    for bad in doc.xpath("//script|//style"):
        bad.getparent().remove(bad)
    out = []
    for tbl in doc.xpath("//table"):
        rows = [[" ".join("".join(td.itertext()).split())
                 for td in tr.xpath("./td|./th")] for tr in tbl.xpath(".//tr")]
        t = md_table(rows)
        if t:
            out.append(t)
    text = "\n".join(ln.strip() for ln in doc.text_content().splitlines() if ln.strip())
    return text + ("\n\n" + "\n\n".join(out) if out else "")


def parse_zip(path):
    with zipfile.ZipFile(path) as z:
        rows = [["파일", "크기(bytes)"]]
        rows += [[i.filename, f"{i.file_size:,}"] for i in z.infolist() if not i.is_dir()]
    return "압축 파일 내용 목록 (본문 미추출):\n\n" + md_table(rows)


HANDLERS = {
    ".hwp": parse_hwp, ".hwpx": parse_hwpx,
    ".pptx": parse_pptx, ".xlsx": parse_xlsx, ".docx": parse_docx,
    ".pdf": parse_pdf, ".csv": parse_csv,
    ".html": parse_html, ".htm": parse_html, ".zip": parse_zip,
}


def main():
    src, out = Path(sys.argv[1]).resolve(), Path(sys.argv[2]).resolve()
    manifest = []
    files = sorted(p for p in src.rglob("*")
                   if p.is_file() and not p.name.startswith("."))

    for f in files:
        rel = f.relative_to(src)
        target = out / rel.with_suffix(rel.suffix + ".md")
        target.parent.mkdir(parents=True, exist_ok=True)
        ext = f.suffix.lower()
        handler = HANDLERS.get(ext)
        status, body = "ok", ""
        try:
            if handler is None:
                status = "skip"
                body = f"_(파싱 대상 아님 — `{ext}` 바이너리. 원본 참조)_"
            else:
                body = (handler(f) or "").strip()
                if not body:
                    status, body = "empty", "_(추출된 텍스트 없음)_"
                elif body.startswith("_(텍스트 레이어 없음"):
                    status = "scan"
        except Exception as e:
            status = "error"
            body = f"_(파싱 실패: {type(e).__name__}: {e})_"

        if len(body) > MAX_CHARS:
            body = body[:MAX_CHARS] + f"\n\n_(이하 생략 — 원문 {len(body):,}자)_"

        stat = f.stat()
        header = "\n".join([
            f"# {f.name}", "",
            f"- 원본 경로: `{rel}`",
            f"- 형식: `{ext}`  |  크기: {stat.st_size:,} bytes",
            f"- 수정일: {datetime.fromtimestamp(stat.st_mtime):%Y-%m-%d %H:%M}",
            f"- 파싱 상태: `{status}`",
            "", "---", "",
        ])
        target.write_text(header + body + "\n", encoding="utf-8")
        manifest.append({"src": str(rel), "parsed": str(target.relative_to(out)),
                         "ext": ext, "bytes": stat.st_size,
                         "mtime": f"{datetime.fromtimestamp(stat.st_mtime):%Y-%m-%d}",
                         "status": status, "chars": len(body)})
        print(f"{status:6s} {rel}", flush=True)

    (out / "_manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=1), encoding="utf-8")
    tally = {}
    for m in manifest:
        tally[m["status"]] = tally.get(m["status"], 0) + 1
    print(f"\n총 {len(manifest)}개 — " + ", ".join(f"{k}={v}" for k, v in sorted(tally.items())))


if __name__ == "__main__":
    main()
