#!/usr/bin/env python3
"""Urban Data Lab native-PPTX toolkit.

Builds fully editable PowerPoint slides from the design.md token set — real
shapes, real text runs, real tables. No rasterised slide images.

Coordinate model: author in the design system's 1600x900 px canvas; this module
converts to EMU/points. A 13.333in slide is 960pt wide, so 1px = 0.6pt = 7620 EMU.

Usage as a library:
    from udl_pptx import Deck, T
    d = Deck()
    s = d.slide(theme="dark")
    s.text(78, 126, 1400, 24, "DATA COLLECTION", T.EYEBROW, color=T.BLUE)
    d.save("out.pptx")
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

PX_TO_EMU = 7620          # 12192000 EMU / 1600 px
PX_TO_PT = 0.6            # 960 pt / 1600 px
CANVAS_W, CANVAS_H = 1600, 900


def px(v):
    return Emu(int(round(v * PX_TO_EMU)))


def pt(v):
    """Design-system px font size -> PowerPoint points."""
    return Pt(round(v * PX_TO_PT, 1))


def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class T:
    """Design tokens (design.md §3, §4). Weight-specific font families are
    required by PowerPoint — CSS numeric weights do not carry over."""

    # brand
    MINT = "#14DEAF"
    TEAL = "#33C3C6"
    BLUE = "#108CEB"
    PURPLE = "#745CED"
    ACCENT = "#00FFCB"
    DEEP = "#0A082B"
    MED_NAVY = "#000050"
    MED_BLUE = "#5080F0"
    MED_PURPLE = "#A050F0"

    # weight -> installed family name
    F = {
        100: "Paperlogy 1 Thin",
        200: "Paperlogy 2 ExtraLight",
        300: "Paperlogy 3 Light",
        400: "Paperlogy 4 Regular",
        500: "Paperlogy 5 Medium",
        600: "Paperlogy 6 SemiBold",
        700: "Paperlogy 7 Bold",
        800: "Paperlogy 8 ExtraBold",
        900: "Paperlogy 9 Black",
    }
    F2 = {
        400: "Freesentation 4 Regular",
        500: "Freesentation 5 Medium",
        700: "Freesentation 7 Bold",
    }

    # (size_px, weight, letter_spacing_px) — design.md §4.2
    DISPLAY = (82, 900, -2.9)
    H1 = (60, 800, -1.8)
    H2 = (40, 700, -1.0)
    H3 = (28, 700, -0.4)
    EYEBROW = (15, 800, 1.8)
    BODY_LG = (22, 500, -0.2)
    BODY = (18, 400, -0.2)
    CAPTION = (13, 600, 0.1)
    META = (12, 500, 0.6)
    KPI = (56, 900, -1.7)


THEMES = {
    "dark": {
        "canvas": "#0A082B",
        "surface": "#141235",
        "surface2": "#1B1940",
        "text": "#FFFFFF",
        "text2": "#B9B6D6",
        "text3": "#8A87AC",
        "border": "#2A2A55",
        "grid": "#232149",
    },
    "light": {
        "canvas": "#FFFFFF",
        "surface": "#F5F6FA",
        "surface2": "#EAF3FD",
        "text": "#0A082B",
        "text2": "#4A4869",
        "text3": "#7B799A",
        "border": "#D8DAE6",
        "grid": "#E6E7EF",
    },
}


class Slide:
    def __init__(self, slide, theme):
        self._s = slide
        self.theme = theme
        self.c = THEMES[theme]

    # --- primitives ---------------------------------------------------

    def rect(self, x, y, w, h, fill=None, line=None, line_w=1,
             radius=None, shape=MSO_SHAPE.RECTANGLE):
        if radius is not None:
            shape = MSO_SHAPE.ROUNDED_RECTANGLE
        sp = self._s.shapes.add_shape(shape, px(x), px(y), px(w), px(h))
        if radius is not None:
            # adjustment is a fraction of the shorter side
            sp.adjustments[0] = min(0.5, radius / min(w, h))
        if fill:
            sp.fill.solid()
            sp.fill.fore_color.rgb = rgb(fill)
        else:
            sp.fill.background()
        if line:
            sp.line.color.rgb = rgb(line)
            sp.line.width = Pt(line_w * PX_TO_PT)
        else:
            sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def gradient(self, x, y, w, h, stops, angle=0):
        """stops: list of (position 0..1, hex). angle in degrees, 0 = left→right."""
        sp = self._s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
        sp.fill.gradient()
        sp.fill.gradient_angle = angle
        gs = sp.fill.gradient_stops
        # python-pptx creates 2 stops by default; reuse then extend via XML if needed
        for i, (pos, color) in enumerate(stops[: len(gs)]):
            gs[i].position = pos
            gs[i].color.rgb = rgb(color)
        sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def text(self, x, y, w, h, content, style, color=None, weight=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, secondary=False,
             line_spacing=None):
        """content: str, or list of (text, color_or_None, weight_or_None) runs."""
        size_px, w_default, spacing = style
        weight = weight or w_default
        color = color or self.c["text"]

        box = self._s.shapes.add_textbox(px(x), px(y), px(w), px(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing

        runs = content if isinstance(content, list) else [(content, None, None)]
        for text, rcolor, rweight in runs:
            r = p.add_run()
            r.text = text
            fam = T.F2 if secondary else T.F
            wt = rweight or weight
            r.font.name = fam.get(wt, fam[min(fam, key=lambda k: abs(k - wt))])
            r.font.size = pt(size_px)
            r.font.color.rgb = rgb(rcolor or color)
            # letter spacing is not exposed by python-pptx; write it directly
            r.font._rPr.set("spc", str(int(spacing * PX_TO_PT * 100)))
        return box

    def line(self, x, y, w, color, thickness=1):
        return self.rect(x, y, w, thickness, fill=color)

    @staticmethod
    def _cell_borders(cell, bottom=None, width_px=1):
        """design.md §6.8: horizontal hairlines only, no vertical rules."""
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            for old in tcPr.findall(qn(tag)):
                tcPr.remove(old)
        # order matters in the CT_TableCellProperties schema
        for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            ln = tcPr.makeelement(qn(tag), {})
            if tag == "a:lnB" and bottom:
                ln.set("w", str(int(width_px * PX_TO_PT * 12700)))
                fill = ln.makeelement(qn("a:solidFill"), {})
                clr = fill.makeelement(qn("a:srgbClr"), {"val": bottom.lstrip("#")})
                fill.append(clr)
                ln.append(fill)
            else:
                ln.append(ln.makeelement(qn("a:noFill"), {}))
            tcPr.append(ln)

    def table(self, x, y, w, h, rows, col_widths=None,
              header_color=None, total_col=True):
        n_r, n_c = len(rows), len(rows[0])
        gf = self._s.shapes.add_table(n_r, n_c, px(x), px(y), px(w), px(h))
        tbl = gf.table
        tbl.first_row = False          # suppress PowerPoint's banded styling
        tbl.horz_banding = False

        if col_widths:
            for i, cw in enumerate(col_widths):
                tbl.columns[i].width = px(cw)

        header_color = header_color or T.MINT
        for ri, row in enumerate(rows):
            tbl.rows[ri].height = px(h / n_r)
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = ""
                cell.margin_left = cell.margin_right = px(16)
                cell.margin_top = cell.margin_bottom = px(8)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(
                    self.c["surface2"] if ri == 0 else self.c["canvas"])
                self._cell_borders(
                    cell, bottom=None if ri == n_r - 1 else self.c["grid"])

                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
                r = p.add_run()
                r.text = str(val)
                if ri == 0:
                    r.font.name, col = T.F[700], header_color
                elif ci == 0:
                    r.font.name, col = T.F[700], self.c["text"]
                elif total_col and ci == n_c - 1:
                    r.font.name, col = T.F[700], self.c["text"]
                else:
                    r.font.name, col = T.F[400], self.c["text2"]
                r.font.size = pt(18)
                r.font.color.rgb = rgb(col)
        return tbl


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = px(CANVAS_W)
        self.prs.slide_height = px(CANVAS_H)

    def slide(self, theme="dark"):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        sl = Slide(s, theme)
        sl.rect(0, 0, CANVAS_W, CANVAS_H, fill=sl.c["canvas"])   # canvas
        return sl

    def save(self, path):
        self.prs.save(path)
        return path
