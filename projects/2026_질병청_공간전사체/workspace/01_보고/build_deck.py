#!/usr/bin/env python3
"""Build the 2026 KDCA spatial-transcriptomics mid-term deck.

Usage:
    python3 build_deck.py content.json 중간보고_초안_v1.pptx

The builder consumes the DECK_SPEC.md JSON schema and emits a fully editable,
native PowerPoint deck. Raster objects are added only for the supplied visual
assets; captions, frames, tables, labels, and all other content remain native.
"""

from __future__ import annotations

import json
import re
import sys
from collections.abc import Iterable
from importlib import import_module
from pathlib import Path
from typing import Any

from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
REPO_ROOT = next(
    parent for parent in HERE.parents if (parent / "scripts" / "udl_pptx.py").is_file()
)
sys.path.insert(0, str(REPO_ROOT / "scripts"))

_udl_pptx = import_module("udl_pptx")
Deck = _udl_pptx.Deck
Slide = _udl_pptx.Slide
T = _udl_pptx.T
pt = _udl_pptx.pt
px = _udl_pptx.px


THEME = "light"
SAFE_X = 78
SAFE_RIGHT = 70
SAFE_W = 1600 - SAFE_X - SAFE_RIGHT
GRID_GAP = 24
COL_W = 99
TITLE_Y = 126
CONTENT_Y = 286
FOOTER_Y = 816

STYLE_TITLE_COMPACT = (34, 700, -0.8)
STYLE_BODY_COMPACT = (16, 400, -0.1)
STYLE_CAPTION_COMPACT = (11, 600, 0.1)
STYLE_KPI_COMPACT = (46, 900, -1.3)

SUPPORTED_TYPES = {
    "cover",
    "agenda",
    "kpi_table",
    "timeline",
    "image_text",
    "image_grid",
    "table",
    "list",
    "closing",
}


def grid_span(columns: int) -> int:
    """Width of N columns in the 12-column design grid."""
    return COL_W * columns + GRID_GAP * (columns - 1)


def fit_title_style(title: str):
    return STYLE_TITLE_COMPACT if len(title) > 48 else T.H2


def resolve_asset(content_dir: Path, relative_path: str) -> Path:
    """Resolve a content asset and select an explicit corrected sibling if present."""
    path = (content_dir / relative_path).resolve()
    corrected = path.with_name(f"{path.stem}_보정{path.suffix}")
    if corrected.is_file():
        path = corrected
    if not path.is_file():
        raise FileNotFoundError(f"이미지 자산을 찾을 수 없습니다: {relative_path}")
    return path


def validate_content(data: dict[str, Any]) -> None:
    if not isinstance(data.get("meta"), dict):
        raise TypeError("content.json에 meta 객체가 필요합니다.")
    slides = data.get("slides")
    if not isinstance(slides, list) or not slides:
        raise ValueError("content.json에 비어 있지 않은 slides 배열이 필요합니다.")

    seen: set[int] = set()
    for index, item in enumerate(slides, start=1):
        no = item.get("no")
        slide_type = item.get("type")
        if not isinstance(no, int) or no <= 0:
            raise ValueError(f"slides[{index - 1}].no는 양의 정수여야 합니다.")
        if no in seen:
            raise ValueError(f"중복 슬라이드 번호: {no}")
        seen.add(no)
        if slide_type not in SUPPORTED_TYPES:
            raise ValueError(f"지원하지 않는 슬라이드 유형: {slide_type!r}")
        if not str(item.get("title", "")).strip():
            raise ValueError(f"{no}번 슬라이드 title이 비어 있습니다.")

        kpis = item.get("kpis", [])
        if len(kpis) > 4:
            raise ValueError(f"{no}번 슬라이드 KPI는 최대 4개입니다.")
        if sum(bool(kpi.get("watch")) for kpi in kpis) > 1:
            raise ValueError(f"{no}번 슬라이드 watch KPI는 최대 1개입니다.")
        bullets = item.get("bullets", [])
        if len(bullets) > 5:
            raise ValueError(f"{no}번 슬라이드 bullets는 최대 5개입니다.")
        if slide_type == "timeline" and len(item.get("events", [])) != 6:
            raise ValueError(f"{no}번 타임라인은 월별 6칸이어야 합니다.")


def add_text(
    s: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    style=T.BODY,
    *,
    color: str | None = None,
    weight: int | None = None,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    secondary: bool = False,
):
    mixed_latin = (
        isinstance(text, str)
        and re.search(r"[A-Za-z]", text)
        and re.search(r"[가-힣]", text)
    )
    content = text
    if mixed_latin:
        parts = re.split(
            r"([A-Za-z][A-Za-z0-9&./+%()~-]*(?: [A-Za-z0-9&./+%()~-]+)*)",
            text,
        )
        content = [(part, None, None) for part in parts if part]

    box = s.text(
        x,
        y,
        w,
        h,
        content,
        style,
        color=color,
        weight=weight,
        align=align,
        anchor=anchor,
        secondary=secondary,
    )
    # python-pptx writes only a:latin for font.name. Mixed Korean/Latin runs
    # need a matching East Asian mapping for PowerPoint and LibreOffice.
    for paragraph in box.text_frame.paragraphs:
        for run in paragraph.runs:
            is_latin = bool(re.search(r"[A-Za-z]", run.text))
            if not secondary and is_latin:
                # Freesentation 7 Bold intermittently drops mixed-script Latin
                # glyphs in LibreOffice; Medium is stable in both renderers.
                run.font.name = T.F2[500]
            typeface = run.font.name
            if not typeface:
                continue
            r_pr = run.font._rPr
            r_pr.set("lang", "en-US" if is_latin else "ko-KR")
            if is_latin:
                continue
            for tag in ("a:ea", "a:cs"):
                node = r_pr.find(qn(tag))
                if node is None:
                    node = r_pr.makeelement(qn(tag), {})
                    r_pr.append(node)
                node.set("typeface", typeface)
    return box


def add_footer(s: Slide, slide_no: int) -> None:
    s.line(SAFE_X, FOOTER_Y, SAFE_W, s.c["border"], thickness=1)
    add_text(
        s,
        1470,
        FOOTER_Y + 10,
        60,
        24,
        f"{slide_no:02d}",
        T.META,
        color=s.c["text3"],
        align=PP_ALIGN.RIGHT,
        secondary=True,
    )


def add_header(s: Slide, item: dict[str, Any], accent_enabled: bool = True) -> None:
    eyebrow_color = s.c["accent"] if accent_enabled else s.c["text3"]
    add_text(
        s,
        SAFE_X,
        54,
        820,
        22,
        item.get("eyebrow", ""),
        T.EYEBROW,
        color=eyebrow_color,
    )
    add_text(
        s,
        1240,
        54,
        290,
        22,
        "2026 · MID-TERM REPORT",
        T.META,
        color=s.c["text3"],
        align=PP_ALIGN.RIGHT,
        secondary=True,
    )
    add_text(
        s,
        SAFE_X,
        TITLE_Y,
        SAFE_W,
        110,
        item["title"],
        fit_title_style(item["title"]),
        color=s.c["text"],
    )


def add_bullets(
    s: Slide,
    bullets: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    numbered: bool = False,
    compact: bool = False,
) -> None:
    values = list(bullets)
    if not values:
        return
    gap = 14
    item_h = (h - gap * (len(values) - 1)) / len(values)
    style = STYLE_BODY_COMPACT if compact else T.BODY
    for idx, text in enumerate(values, start=1):
        top = y + (idx - 1) * (item_h + gap)
        marker = f"{idx:02d}" if numbered else "•"
        marker_w = 42 if numbered else 24
        add_text(
            s,
            x,
            top + 2,
            marker_w,
            min(30, item_h),
            marker,
            T.CAPTION,
            color=s.c["text3"],
            secondary=numbered,
        )
        add_text(
            s,
            x + marker_w,
            top,
            w - marker_w,
            item_h,
            text,
            style,
            color=s.c["text2"],
        )


def _set_picture_rounding(picture, radius_px: float, width_px: float, height_px: float) -> None:
    """Change the picture geometry to a rounded rectangle."""
    sp_pr = picture._pic.spPr
    geom = sp_pr.find(qn("a:prstGeom"))
    if geom is None:
        geom = sp_pr.makeelement(qn("a:prstGeom"), {"prst": "roundRect"})
        sp_pr.append(geom)
    else:
        geom.set("prst", "roundRect")

    av_lst = geom.find(qn("a:avLst"))
    if av_lst is None:
        av_lst = geom.makeelement(qn("a:avLst"), {})
        geom.append(av_lst)
    for old in list(av_lst):
        av_lst.remove(old)
    adjustment = int(min(50000, max(0, radius_px / min(width_px, height_px) * 100000)))
    av_lst.append(
        av_lst.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {adjustment}"})
    )


def _contain_size(image_path: Path, frame_w: float, frame_h: float) -> tuple[float, float]:
    """Return dimensions that fit the complete image inside the frame."""
    with Image.open(image_path) as image:
        source_w, source_h = image.size
    source_ratio = source_w / source_h
    frame_ratio = frame_w / frame_h
    if source_ratio > frame_ratio:
        return frame_w, frame_w / source_ratio
    return frame_h * source_ratio, frame_h


def add_picture_contain(
    s: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    path: Path,
    *,
    description: str = "",
    radius: float | None = None,
) -> Any:
    """Add one uncropped PICTURE, centered in its frame at the original ratio."""
    picture_w, picture_h = _contain_size(path, w, h)
    picture_x = x + (w - picture_w) / 2
    picture_y = y + (h - picture_h) / 2
    picture = s._s.shapes.add_picture(
        str(path),
        px(picture_x),
        px(picture_y),
        px(picture_w),
        px(picture_h),
    )
    if radius is not None:
        _set_picture_rounding(picture, radius, picture_w, picture_h)
    picture.name = f"PHOTO · {path.name}"
    picture._pic.nvPicPr.cNvPr.set("descr", description or path.name)
    return picture


def image_card(
    s: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    path: Path,
    caption: str,
) -> Any:
    """Add one uncropped, contained PICTURE plus native frame and caption."""
    caption_h = 48
    s.rect(x, y, w, h, fill=s.c["surface"], line=s.c["border"], radius=14)
    picture = add_picture_contain(
        s,
        x,
        y,
        w,
        h - caption_h,
        path,
        description=caption,
        radius=14,
    )

    s.rect(x, y + h - caption_h, w, caption_h, fill=s.c["canvas"])
    s.rect(x, y + h - caption_h, 5, caption_h, fill=s.c["accent"])
    add_text(
        s,
        x + 18,
        y + h - caption_h + 10,
        w - 32,
        caption_h - 16,
        caption,
        T.CAPTION,
        color=s.c["text2"],
    )
    return picture


def add_kpi_tiles(
    s: Slide,
    kpis: list[dict[str, Any]],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    if not kpis:
        return
    gap = GRID_GAP
    tile_w = (w - gap * (len(kpis) - 1)) / len(kpis)
    for idx, kpi in enumerate(kpis):
        left = x + idx * (tile_w + gap)
        watch = bool(kpi.get("watch"))
        line = s.c["accent"] if watch else s.c["border"]
        s.rect(left, y, tile_w, h, fill=s.c["surface"], line=line, line_w=2 if watch else 1, radius=14)
        add_text(
            s,
            left + 20,
            y + 18,
            tile_w - 40,
            26,
            str(kpi.get("label", "")),
            T.CAPTION,
            color=s.c["text2"],
        )
        value = str(kpi.get("value", ""))
        unit = str(kpi.get("unit", ""))
        value_box = add_text(
            s,
            left + 20,
            y + 48,
            tile_w - 40,
            44,
            [(value, s.c["accent"] if watch else s.c["text"], 900), (f" {unit}", s.c["text2"], 600)],
            STYLE_KPI_COMPACT,
        )
        if len(value_box.text_frame.paragraphs[0].runs) > 1:
            value_box.text_frame.paragraphs[0].runs[1].font.size = pt(22)
        add_text(
            s,
            left + 20,
            y + h - 25,
            tile_w - 40,
            16,
            str(kpi.get("note", "")),
            STYLE_CAPTION_COMPACT,
            color=s.c["text3"],
        )


def add_table_from_spec(
    s: Slide,
    table_spec: dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    head = table_spec.get("head", [])
    rows = table_spec.get("rows", [])
    if not head:
        return
    matrix = [head, *rows]
    col_count = len(head)
    first = 235 if col_count >= 5 else 290
    remaining = (w - first) / (col_count - 1)
    widths = [first, *([remaining] * (col_count - 1))]
    table = s.table(
        x,
        y,
        w,
        h,
        matrix,
        col_widths=widths,
        header_color=s.c["accent"],
    )
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    typeface = run.font.name
                    if not typeface:
                        continue
                    r_pr = run.font._rPr
                    r_pr.set("lang", "ko-KR")
                    for tag in ("a:ea", "a:cs"):
                        node = r_pr.find(qn(tag))
                        if node is None:
                            node = r_pr.makeelement(qn(tag), {})
                            r_pr.append(node)
                        node.set("typeface", typeface)



def fit_card_height(texts, card_w, font_px=18, top=40, bottom=18, indent=72,
                    min_h=96, max_h=None):
    """Height a card needs to hug its text.

    Each text is its own card, so the tallest one governs — summing them
    over-sizes every card and leaves a large empty block underneath.
    Korean glyphs are ~1em wide, Latin ~0.55em.
    """
    usable = max(1.0, card_w - indent - 36)
    worst = 1
    for t in texts:
        width = sum(1.0 if ord(ch) > 0x2E80 else 0.55 for ch in str(t))
        worst = max(worst, max(1, int(width * font_px / usable + 0.999)))
    h = top + worst * font_px * 1.6 + bottom
    if max_h:
        h = min(h, max_h)
    return max(min_h, round(h))


def add_insight_cards(
    s: Slide,
    bullets: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    if not bullets:
        return
    gap = GRID_GAP
    card_w = (w - gap * (len(bullets) - 1)) / len(bullets)
    h = min(h, fit_card_height(bullets, card_w, font_px=17, top=44, bottom=18, indent=36, min_h=96, max_h=h))
    for idx, bullet in enumerate(bullets):
        left = x + idx * (card_w + gap)
        s.rect(left, y, card_w, h, fill=s.c["surface"], line=s.c["border"], radius=14)
        add_text(
            s,
            left + 18,
            y + 14,
            32,
            22,
            f"{idx + 1:02d}",
            T.META,
            color=s.c["text3"],
            secondary=True,
        )
        add_text(
            s,
            left + 18,
            y + 40,
            card_w - 36,
            h - 50,
            bullet,
            STYLE_BODY_COMPACT,
            color=s.c["text2"],
        )


def render_cover(
    deck: Deck,
    item: dict[str, Any],
    meta: dict[str, Any],
    content_dir: Path,
) -> None:
    s = deck.slide(theme=THEME)
    s.rect(0, 0, 10, 900, fill=s.c["accent"])
    add_text(s, SAFE_X, 70, 650, 24, item.get("eyebrow", ""), T.EYEBROW, color=s.c["accent"])
    add_text(s, SAFE_X, 176, 660, 270, item["title"], T.H1, color=s.c["text"])
    add_text(
        s,
        SAFE_X,
        500,
        650,
        72,
        item.get("subtitle", meta.get("subtitle", "")),
        T.BODY_LG,
        color=s.c["text2"],
    )
    s.line(SAFE_X, 576, 560, s.c["border"], thickness=1)
    add_text(s, SAFE_X, 600, 650, 28, item.get("org", meta.get("org", "")), T.BODY, color=s.c["text"])
    add_text(
        s,
        SAFE_X,
        636,
        650,
        26,
        item.get("partners", meta.get("partners", "")),
        T.CAPTION,
        color=s.c["text2"],
    )
    add_text(
        s,
        SAFE_X,
        674,
        300,
        22,
        item.get("date", meta.get("date", "")),
        T.META,
        color=s.c["text3"],
        secondary=True,
    )

    hero_path = resolve_asset(content_dir, "자산/cover/spatial_xenium.jpg")
    hero_x, hero_y, hero_w, hero_h = 790, 122, 740, 520
    s.rect(
        hero_x,
        hero_y,
        hero_w,
        hero_h,
        fill=s.c["surface"],
        line=s.c["border"],
        radius=14,
    )
    add_picture_contain(
        s,
        hero_x + 16,
        hero_y + 16,
        hero_w - 32,
        hero_h - 32,
        hero_path,
        description="Xenium 공간전사체 분석 이미지",
        radius=14,
    )

    logo_specs = [
        ("자산/ci_final/urbandatalab.png", 54),
        ("자산/ci_final/gil.png", 38),
        ("자산/ci_final/portrai.png", 42),
        ("자산/ci_final/nih.png", 38),
    ]
    logo_y = 732
    logo_w = 280
    logo_gap = 70
    for idx, (relative_path, visual_h) in enumerate(logo_specs):
        logo_x = SAFE_X + idx * (logo_w + logo_gap)
        logo_path = resolve_asset(content_dir, relative_path)
        add_picture_contain(
            s,
            logo_x,
            logo_y + (58 - visual_h) / 2,
            logo_w,
            visual_h,
            logo_path,
            description=logo_path.stem,
        )
    add_footer(s, item["no"])


def render_agenda(deck: Deck, item: dict[str, Any]) -> None:
    s = deck.slide(theme=THEME)
    add_header(s, item)
    items = item.get("items", [])
    for idx, label in enumerate(items[:8]):
        col = 0 if idx < 4 else 1
        row = idx if idx < 4 else idx - 4
        left = SAFE_X + col * (grid_span(6) + GRID_GAP)
        top = CONTENT_Y + row * 112
        s.rect(left, top, grid_span(6), 88, fill=s.c["surface"], line=s.c["border"], radius=14)
        add_text(
            s,
            left + 20,
            top + 20,
            52,
            32,
            f"{idx + 1:02d}",
            T.H3,
            color=s.c["accent"],
            secondary=True,
        )
        clean = label.split(". ", 1)[-1]
        add_text(s, left + 88, top + 21, grid_span(6) - 108, 46, clean, T.BODY, color=s.c["text"])
    add_footer(s, item["no"])


def render_kpi_table(deck: Deck, item: dict[str, Any]) -> None:
    s = deck.slide(theme=THEME)
    has_watch = any(bool(kpi.get("watch")) for kpi in item.get("kpis", []))
    add_header(s, item, accent_enabled=not has_watch)
    add_kpi_tiles(s, item.get("kpis", []), SAFE_X, 244, SAFE_W, 128)
    add_table_from_spec(s, item.get("table", {}), SAFE_X, 396, SAFE_W, 272)
    add_insight_cards(s, item.get("bullets", [])[:3], SAFE_X, 692, SAFE_W, 100)
    add_footer(s, item["no"])


def timeline(
    s: Slide,
    events: list[dict[str, Any]],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    gap = GRID_GAP
    card_w = (w - gap * 5) / 6
    # hug the content: tallest title + desc governs, not a fixed block
    title_h = fit_card_height([e.get("title", "") for e in events], card_w,
                              font_px=20, top=0, bottom=0, indent=36, min_h=52)
    desc_h = fit_card_height([e.get("desc", "") for e in events], card_w,
                             font_px=15, top=0, bottom=0, indent=36, min_h=60)
    h = min(h, 84 + 24 + title_h + 16 + desc_h + 24)
    line_y = y + 50
    s.line(x, line_y, w, s.c["border"], thickness=2)
    for idx, event in enumerate(events):
        left = x + idx * (card_w + gap)
        center = left + card_w / 2
        s.rect(center - 7, line_y - 7, 14, 14, fill=s.c["accent"], radius=7)
        add_text(
            s,
            left,
            y,
            card_w,
            26,
            str(event.get("month", "")),
            T.CAPTION,
            color=s.c["accent"],
            align=PP_ALIGN.CENTER,
        )
        s.rect(left, y + 84, card_w, h - 84, fill=s.c["surface"], line=s.c["border"], radius=14)
        add_text(
            s,
            left + 18,
            y + 108,
            card_w - 36,
            title_h,
            str(event.get("title", "")),
            (20, 700, -0.2),
            color=s.c["text"],
        )
        add_text(
            s,
            left + 18,
            y + 108 + title_h + 16,
            card_w - 36,
            desc_h,
            str(event.get("desc", "")),
            STYLE_BODY_COMPACT,
            color=s.c["text2"],
        )


def render_timeline(deck: Deck, item: dict[str, Any]) -> None:
    s = deck.slide(theme=THEME)
    add_header(s, item)
    timeline(s, item.get("events", []), SAFE_X, 260, SAFE_W, 520)
    add_footer(s, item["no"])


def image_items(item: dict[str, Any]) -> list[dict[str, str]]:
    if item.get("images"):
        return [
            {"path": str(value["path"]), "caption": str(value.get("caption", ""))}
            for value in item["images"]
        ]
    if item.get("image"):
        return [
            {
                "path": str(item["image"]),
                "caption": str(item.get("image_caption", "")),
            }
        ]
    return []


def render_image_text(deck: Deck, item: dict[str, Any], content_dir: Path) -> None:
    s = deck.slide(theme=THEME)
    add_header(s, item)
    visuals = image_items(item)
    image_w = grid_span(7)
    text_x = SAFE_X + image_w + GRID_GAP
    text_w = grid_span(5)

    if len(visuals) == 1:
        visual = visuals[0]
        image_card(
            s,
            SAFE_X,
            CONTENT_Y,
            image_w,
            496,
            resolve_asset(content_dir, visual["path"]),
            visual["caption"],
        )
    else:
        gap = 16
        card_w = (image_w - gap * (len(visuals) - 1)) / len(visuals)
        for idx, visual in enumerate(visuals):
            image_card(
                s,
                SAFE_X + idx * (card_w + gap),
                CONTENT_Y,
                card_w,
                496,
                resolve_asset(content_dir, visual["path"]),
                visual["caption"],
            )

    s.rect(text_x, CONTENT_Y, text_w, 496, fill=s.c["surface"], line=s.c["border"], radius=14)
    add_text(s, text_x + 24, CONTENT_Y + 24, text_w - 48, 24, "KEY POINTS", T.CAPTION, color=s.c["text3"])
    add_bullets(
        s,
        item.get("bullets", []),
        text_x + 24,
        CONTENT_Y + 72,
        text_w - 48,
        390,
        compact=True,
    )
    add_footer(s, item["no"])


def render_image_grid(deck: Deck, item: dict[str, Any], content_dir: Path) -> None:
    s = deck.slide(theme=THEME)
    add_header(s, item)
    visuals = image_items(item)
    card_w = grid_span(3)
    card_h = 270
    for idx, visual in enumerate(visuals[:4]):
        image_card(
            s,
            SAFE_X + idx * (card_w + GRID_GAP),
            248,
            card_w,
            card_h,
            resolve_asset(content_dir, visual["path"]),
            visual["caption"],
        )
    if item.get("table"):
        add_table_from_spec(s, item["table"], SAFE_X, 540, SAFE_W, 246)
    else:
        add_insight_cards(s, item.get("bullets", [])[:3], SAFE_X, 548, SAFE_W, 230)
    add_footer(s, item["no"])


def render_table(deck: Deck, item: dict[str, Any]) -> None:
    s = deck.slide(theme=THEME)
    add_header(s, item)
    if item.get("kpis"):
        add_kpi_tiles(s, item["kpis"], SAFE_X, 238, SAFE_W, 124)
        table_y = 388
        table_h = 292
    else:
        table_y = CONTENT_Y
        table_h = 380
    add_table_from_spec(s, item.get("table", {}), SAFE_X, table_y, SAFE_W, table_h)
    add_insight_cards(s, item.get("bullets", [])[:3], SAFE_X, 704, SAFE_W, 88)
    add_footer(s, item["no"])


def render_list(deck: Deck, item: dict[str, Any]) -> None:
    s = deck.slide(theme=THEME)
    add_header(s, item)
    bullets = item.get("bullets", [])
    card_w = grid_span(6)
    card_h = fit_card_height(
        [b.split(". ", 1)[-1] for b in bullets[:4]], card_w,
        font_px=18, top=30, bottom=24, indent=120, min_h=132, max_h=240)
    for idx, bullet in enumerate(bullets[:4]):
        col = idx % 2
        row = idx // 2
        left = SAFE_X + col * (card_w + GRID_GAP)
        top = CONTENT_Y + row * (card_h + GRID_GAP)
        s.rect(left, top, card_w, card_h, fill=s.c["surface"], line=s.c["border"], radius=14)
        add_text(
            s,
            left + 24,
            top + 24,
            54,
            34,
            f"{idx + 1:02d}",
            T.H3,
            color=s.c["accent"],
            secondary=True,
        )
        cleaned = bullet.split(". ", 1)[-1]
        add_text(
            s,
            left + 96,
            top + 24,
            card_w - 120,
            card_h - 48,
            cleaned,
            T.BODY,
            color=s.c["text2"],
        )
    add_footer(s, item["no"])


def render_closing(deck: Deck, item: dict[str, Any]) -> None:
    s = deck.slide(theme=THEME)
    s.rect(0, 0, 10, 900, fill=s.c["accent"])
    s.rect(1030, 0, 570, 900, fill=s.c["surface"])
    add_text(s, SAFE_X, 166, 700, 24, item.get("eyebrow", ""), T.EYEBROW, color=s.c["accent"])
    add_text(s, SAFE_X, 286, 760, 100, item["title"], T.DISPLAY, color=s.c["text"])
    add_text(
        s,
        SAFE_X,
        430,
        820,
        100,
        item.get("subtitle", ""),
        T.BODY_LG,
        color=s.c["text2"],
    )
    s.line(SAFE_X, 576, 680, s.c["border"], thickness=1)
    add_text(s, SAFE_X, 610, 820, 34, item.get("org", ""), T.BODY, color=s.c["text"])
    add_text(s, 1100, 380, 360, 48, "URBAN DATA LAB", T.H3, color=s.c["accent"], align=PP_ALIGN.CENTER)
    add_text(
        s,
        1100,
        448,
        360,
        24,
        "BUILD EVIDENCE · DELIVER BETTER CARE",
        T.META,
        color=s.c["text2"],
        align=PP_ALIGN.CENTER,
        secondary=True,
    )
    add_footer(s, item["no"])


def build_deck(data: dict[str, Any], content_path: Path, output_path: Path) -> Path:
    validate_content(data)
    deck = Deck()
    meta = data["meta"]
    content_dir = content_path.parent
    renderers = {
        "cover": lambda item: render_cover(deck, item, meta, content_dir),
        "agenda": lambda item: render_agenda(deck, item),
        "kpi_table": lambda item: render_kpi_table(deck, item),
        "timeline": lambda item: render_timeline(deck, item),
        "image_text": lambda item: render_image_text(deck, item, content_dir),
        "image_grid": lambda item: render_image_grid(deck, item, content_dir),
        "table": lambda item: render_table(deck, item),
        "list": lambda item: render_list(deck, item),
        "closing": lambda item: render_closing(deck, item),
    }
    for item in sorted(data["slides"], key=lambda value: value["no"]):
        renderers[item["type"]](item)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    deck.save(str(output_path))
    return output_path


def main(argv: list[str]) -> int:
    if len(argv) != 3:
        print(
            "사용법: python3 build_deck.py content.json 중간보고_초안_v2.pptx",
            file=sys.stderr,
        )
        return 2

    content_path = Path(argv[1]).expanduser().resolve()
    output_path = Path(argv[2]).expanduser().resolve()
    with content_path.open("r", encoding="utf-8") as handle:
        data = json.load(handle)
    built = build_deck(data, content_path, output_path)

    # Fail loudly if a future change accidentally rasterizes a full slide.
    picture_count = sum(
        1
        for slide in DeckInspection.open(built)
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    expected_pictures = sum(len(image_items(item)) for item in data["slides"])
    expected_pictures += 5 * sum(item["type"] == "cover" for item in data["slides"])
    if picture_count != expected_pictures:
        raise RuntimeError(
            f"PICTURE 개수 불일치: 실제 {picture_count}, 콘텐츠 사진 {expected_pictures}"
        )
    print(f"생성 완료: {built}")
    print(f"슬라이드: {len(data['slides'])}장 · PICTURE: {picture_count}개")
    return 0


class DeckInspection:
    """Late import wrapper keeps builder construction centered on udl_pptx.Deck."""

    @staticmethod
    def open(path: Path):
        from pptx import Presentation

        return Presentation(str(path)).slides


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
