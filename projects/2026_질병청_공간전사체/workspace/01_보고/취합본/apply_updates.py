#!/usr/bin/env python3
"""문승완본(v0.5_20260727_문승완)의 갱신 내용을 취합 기준본에 반영한다.

기준본: 2026_공간바이오마커_중간보고_초안_v2.pptx (19장)
반영본: 공간바이오마커_중간보고회_v0.5_20260727_문승완.pptx (10장)

서식을 보존하기 위해 셀·런 단위로 텍스트만 교체하고, 새 내용은 기존 표의 열/행을 확장해 넣는다.
콘솔 캡처는 삽입하지 않고 수치만 네이티브 표로 옮긴다(design.md §12.3).
"""
import copy
import sys

from pptx import Presentation

SRC = "2026_공간바이오마커_중간보고_초안_v2.pptx"
OUT = "2026_공간바이오마커_중간보고_초안_v3_취합.pptx"

# 길병원 집행 갱신에 따른 총액 재계산
URBAN, PORTRAI, GIL = 20_296_612, 29_936_083, 51_176_435
TOTAL = URBAN + PORTRAI + GIL          # 101,409,130
BUDGET = 300_000_000

# CellVit 최초 학습 결과 (문승완본 s6 · 평가 패치 27,468개 동일)
CELLVIT = {"DICE": "0.6427", "AJI": "0.2537", "PQ": "0.1782",
           "SQ": "0.4827", "DQ": "0.2792"}


def set_cell(cell, text):
    """첫 런의 서식을 유지한 채 셀 텍스트를 교체한다."""
    tf = cell.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.add_run().text = text


def clone_column(table, src_idx):
    """표의 마지막 열을 복제해 새 열을 만든다(서식 유지)."""
    grid = table._tbl.tblGrid
    new_gc = copy.deepcopy(grid[src_idx])
    grid.append(new_gc)
    for tr in table._tbl.tr_lst:
        new_tc = copy.deepcopy(tr.tc_lst[src_idx])
        tr.append(new_tc)


def replace_text_everywhere(slide, pairs):
    n = 0
    for sh in slide.shapes:
        frames = []
        if sh.has_text_frame:
            frames.append(sh.text_frame)
        if sh.has_table:
            frames += [c.text_frame for r in sh.table.rows for c in r.cells]
        for tf in frames:
            for p in tf.paragraphs:
                for r in p.runs:
                    for old, new in pairs:
                        if old in r.text:
                            r.text = r.text.replace(old, new)
                            n += 1
    return n


def main():
    prs = Presentation(SRC)
    log = []

    # ── s15 AI: HoverNet 표에 CellVit 열 추가 ───────────────────────
    s15 = prs.slides[14]
    for sh in s15.shapes:
        if not sh.has_table:
            continue
        t = sh.table
        if t.cell(0, 0).text.strip() != "평가 지표":
            continue
        # '값' -> 'HoverNet', 새 열 'CellVit (최초 학습)'
        set_cell(t.cell(0, 1), "HoverNet")
        clone_column(t, len(t.columns) - 1)
        last = len(t.columns) - 1
        set_cell(t.cell(0, last), "CellVit (최초 학습)")
        for ri in range(1, len(t.rows)):
            metric = t.cell(ri, 0).text.strip()
            set_cell(t.cell(ri, last), CELLVIT.get(metric, "—"))
        # 의미 열을 3번째로 유지하되 폭 재배분
        total_w = sum(c.width for c in t.columns)
        widths = [0.16, 0.18, 0.42, 0.24]
        for i, w in enumerate(widths[: len(t.columns)]):
            t.columns[i].width = int(total_w * w)
        log.append("s15 AI: CellVit 최초 학습 결과 열 추가 (평가 패치 27,468개 동일 기준)")
        break
    replace_text_everywhere(s15, [
        ("Class Imbalance 극복을 위한 Loss 개선 및 CellVit 투트랙 교차 검증 추진",
         "HoverNet 개선 학습을 진행하고 CellVit 최초 학습 결과를 확보하여 투트랙 비교 기반을 마련함"),
    ])

    # ── s17 BUDGET: 길병원 집행 갱신 + 총액 재계산 ──────────────────
    s17 = prs.slides[16]
    n = replace_text_everywhere(s17, [
        ("46,027,242", f"{GIL:,}"),
        ("52.90%", "58.83%"),
        ("52.90", "58.83"),
        ("88,649,625", f"{TOTAL:,}"),
        ("29.50%", f"{TOTAL/BUDGET*100:.2f}%"),
        ("29.5%", f"{TOTAL/BUDGET*100:.1f}%"),
        ("88.6", f"{TOTAL/1_000_000:.1f}"),
        ("46.0백만원", f"{GIL/1_000_000:.1f}백만원"),
        ("2026-05-27", "2026-07-27"),
        ("211,3", f"{(BUDGET-TOTAL)//1000%1000:,}"),
    ])
    log.append(f"s17 예산: 길병원 {GIL:,}원(58.83%) 반영 · 총액 {TOTAL:,}원({TOTAL/BUDGET*100:.1f}%) 재계산 [{n}건 치환]")

    # ── s14 CLINICAL: 임상 매칭 진척 반영 ───────────────────────────
    s14 = prs.slides[13]
    n = replace_text_everywhere(s14, [
        ("892건 임상 데이터 8월부터 실측 매칭 수행",
         "Xenium 실시 50명 우선 매칭 후 잔여분 순차 매칭"),
        ("7월부터 매칭 예정", "Xenium 실시 50명 우선 매칭 진행"),
    ])
    log.append(f"s14 임상: Xenium 50명 우선 매칭 반영 [{n}건]")

    # ── s18 NEXT STEPS: SOP·cell mask 추가 ─────────────────────────
    s18 = prs.slides[17]
    n = replace_text_everywhere(s18, [
        ("AI 모델 투트랙 성능 평가: HoverNet과 CellVit 모델 간 비교 평가를 통해 병리-공간전사체 연계 성능 확보.",
         "AI 모델 투트랙 성능 평가: CellVit 파라미터 튜닝으로 HoverNet 대비 성능 개선을 추진하고, "
         "SOP 절차서 양식 수령에 따라 절차서를 선제 준비."),
        ("7월 중 실험 완료 후 8월 내 질병청 완제 기탁 완료.",
         "잔여 블록 4개 수령분까지 생산 100% 완료 후 cell mask 생성 및 8월 내 질병청 완제 기탁."),
    ])
    log.append(f"s18 계획: SOP 절차서·cell mask 생성 반영 [{n}건]")

    prs.save(OUT)
    print(f"저장: {OUT}  ({len(prs.slides)}장)")
    for x in log:
        print("  ·", x)


if __name__ == "__main__":
    sys.exit(main())
