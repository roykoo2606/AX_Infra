#!/usr/bin/env python3
"""QA gate for the interim-report deck — DECK_SPEC.md 검수 기준을 기계적으로 검사한다.

Usage: python3 qa_deck.py 중간보고_초안_v1.pptx
"""
import sys
from collections import Counter

from pptx import Presentation
from pptx.util import Emu

FORBIDDEN_FONTS = ("AppleMyungjo", "Songti SC", "Apple SD Gothic Neo", "Pretendard")
FORBIDDEN_COLORS = ("E0B357", "B0862E", "E2B24C", "00FFCB")

# DECK_SPEC.md "확정 사실" — 등장하면 안 되는 값
BANNED_CLAIMS = {
    "총 54건": "공간전사체 총계 단정 (원본 계산 불일치)",
    "총 64건": "공간전사체 총계 단정",
    "발주자 목표": "6,000건은 자체 목표",
    "발주처 목표": "6,000건은 자체 목표",
}
REQUIRED = {
    "2026-07-27": "사업비 기준일",
    "27.3": "전문의 리뷰 진척률 (v0.5 기준)",
    "0.00%": "임상 연계 미착수",
}
# 대외 보고본에는 내부 시스템 실측치와 내부 문서 참조를 넣지 않는다
BANNED_INTERNAL = ("7,919", "7,898", "2,582", "5,353", "v0.5",
                   # 갱신 전 수치 — 최신본에 남아 있으면 표와 KPI가 어긋난다
                   "46,027,242", "88,649,625", "52.90", "29.50")


def all_text(prs):
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
            if sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells:
                        out.append(c.text)
    return "\n".join(out)


def fonts_used(prs):
    f = Counter()
    for s in prs.slides:
        for sh in s.shapes:
            frames = []
            if sh.has_text_frame:
                frames.append(sh.text_frame)
            if sh.has_table:
                frames += [c.text_frame for r in sh.table.rows for c in r.cells]
            for tf in frames:
                for p in tf.paragraphs:
                    for run in p.runs:
                        if run.font.name:
                            f[run.font.name] += 1
    return f


def main(path):
    prs = Presentation(path)
    text = all_text(prs)
    fails, warns = [], []

    print(f"슬라이드 {len(prs.slides)}장  ·  "
          f"{Emu(prs.slide_width).inches:.3f} × {Emu(prs.slide_height).inches:.3f} in")
    if abs(Emu(prs.slide_width).inches - 13.333) > 0.01:
        fails.append("슬라이드 폭이 13.333in이 아님")

    # 0) 콘솔 캡처 — 기본은 네이티브 표로 대체하나, 원자료를 그대로 싣기로 한 경우
    #    (2026-07-29 결정: 문승완 자료 3장은 이미지 전량 사용)에는 경고만 남긴다.
    CONSOLE_SIZES = {(1709, 247), (1720, 650)}
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.shape_type == 13 and sh.image.size in CONSOLE_SIZES:
                warns.append(f"슬라이드 {i}: 콘솔 캡처 이미지 — 의도한 것인지 확인")

    # 1) 네이티브 여부 — 사진은 허용, 슬라이드 전면 래스터는 금지
    print("\n[네이티브 검사]")
    for i, s in enumerate(prs.slides, 1):
        pics = [sh for sh in s.shapes if sh.shape_type == 13]
        full = [p for p in pics
                if p.width >= prs.slide_width * 0.95 and p.height >= prs.slide_height * 0.95]
        others = [sh for sh in s.shapes if sh.shape_type != 13]
        if full:
            fails.append(f"슬라이드 {i}: 전면 래스터 이미지 {len(full)}개")
        print(f"  s{i:02d}  사진 {len(pics)}  네이티브 개체 {len(others)}")
        if not others:
            fails.append(f"슬라이드 {i}: 네이티브 개체가 없음")

    # 2) 폰트
    print("\n[폰트]")
    fu = fonts_used(prs)
    for name, n in fu.most_common():
        flag = " ❌" if any(b in name for b in FORBIDDEN_FONTS) else ""
        print(f"  {n:5d}  {name}{flag}")
        if flag:
            fails.append(f"금지 폰트 사용: {name}")
    if not any(n.startswith("Paperlogy") for n in fu):
        fails.append("Paperlogy 미사용")

    # 3) 금지 색
    print("\n[금지 색]")
    xml = "".join(s._element.xml for s in prs.slides)
    for c in FORBIDDEN_COLORS:
        if c in xml.upper():
            fails.append(f"금지 색 {c} 사용")
    print("  없음" if not any(c in xml.upper() for c in FORBIDDEN_COLORS) else "  위 실패 참조")

    # 4) 사실 규칙
    print("\n[사실 규칙]")
    for k, why in BANNED_CLAIMS.items():
        if k in text:
            fails.append(f"금지 표현 '{k}' — {why}")
    for k in BANNED_INTERNAL:
        if k in text:
            fails.append(f"내부 전용 값·참조 '{k}' 노출 — 대외 보고본에서 제거할 것")
    for k, why in REQUIRED.items():
        if k not in text:
            warns.append(f"필수 값 '{k}' 미발견 — {why}")
    print("  금지 표현 없음" if not any(k in text for k in BANNED_CLAIMS) else "  위 실패 참조")

    print("\n" + "=" * 56)
    if fails:
        print(f"실패 {len(fails)}건")
        for f in fails:
            print("  ❌", f)
    else:
        print("실패 없음 ✅")
    for w in warns:
        print("  ⚠️ ", w)
    return 1 if fails else 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1]))
